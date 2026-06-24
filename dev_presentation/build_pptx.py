"""Build a high-quality .pptx from the deck.

Each slide is the pixel-perfect rendered HTML slide (so no design is lost), the
deck's speaker notes are attached, and the two demo clips are embedded on the Demo
slide as playable audio (PowerPoint plays embedded audio, unlike PDF).

Run:  python build_pptx.py
Output: ~/Downloads/Project13_presentation.pptx
"""

from __future__ import annotations

import re
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
STANDALONE = HERE / 'presentation_standalone.html'
SRC = HERE / 'presentation.html'
OUT = Path.home() / 'Downloads' / 'Project13_presentation.pptx'


def _find_chrome() -> str | None:
    """Locate a Playwright headless-shell binary in the current user's cache.

    Falls back to None (Playwright's bundled browser) so the build is portable
    across machines instead of relying on one hard-coded absolute path.
    """
    base = Path.home() / 'Library' / 'Caches' / 'ms-playwright'
    hits = sorted(
        base.glob(
            'chromium_headless_shell-*/chrome-headless-shell-*/chrome-headless-shell'
        )
    )
    return str(hits[-1]) if hits else None


CHROME = _find_chrome()
TMP = Path(tempfile.mkdtemp(prefix='pptx_shots_'))
DEMO_AUDIO = [
    HERE / 'assets/audio/demo_base.wav',
    HERE / 'assets/audio/demo_finetuned.wav',
]


def notes_per_slide() -> list[str]:
    """Pull each section's speaker notes from the source HTML, in order."""
    html = SRC.read_text(encoding='utf-8')
    sections = re.findall(r'<section\b.*?</section>', html, re.S)
    out = []
    for s in sections:
        m = re.search(r'<aside class="notes">(.*?)</aside>', s, re.S)
        out.append(re.sub(r'<[^>]+>', '', m.group(1)).strip() if m else '')
    return out


def render_slides() -> list[Path]:
    TMP.mkdir(parents=True, exist_ok=True)
    url = 'file://' + str(STANDALONE.resolve())
    shots = []
    with sync_playwright() as p:
        b = p.chromium.launch(executable_path=CHROME)
        pg = b.new_page(viewport={'width': 1280, 'height': 720}, device_scale_factor=2)
        pg.goto(url)
        pg.wait_for_timeout(1200)
        pg.wait_for_function("typeof Reveal!=='undefined' && Reveal.isReady()")
        n = pg.evaluate('Reveal.getTotalSlides()')
        for i in range(n):
            pg.evaluate(f'Reveal.slide({i})')
            pg.wait_for_timeout(350)
            f = TMP / f'{i:02d}.png'
            pg.screenshot(path=str(f))
            shots.append(f)
        b.close()
    return shots


def main() -> None:
    shots = render_slides()
    notes = notes_per_slide()
    print(f'rendered {len(shots)} slides; {sum(1 for x in notes if x)} have notes')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # which rendered slide is the Demo (to embed audio)?
    html = SRC.read_text(encoding='utf-8')
    titles = [
        re.sub(
            r'<[^>]+>',
            '',
            (
                re.search(r'<h[12][^>]*>(.*?)</h[12]>', s, re.S) or re.match('', '')
            ).group(1),
        ).strip()
        if re.search(r'<h[12][^>]*>(.*?)</h[12]>', s, re.S)
        else ''
        for s in re.findall(r'<section\b.*?</section>', html, re.S)
    ]
    demo_idx = next((i for i, t in enumerate(titles) if t == 'Demo'), None)

    audio_embedded = 0
    for i, img in enumerate(shots):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        if notes[i] if i < len(notes) else '':
            slide.notes_slide.notes_text_frame.text = notes[i]
        if i == demo_idx:
            for k, wav in enumerate(DEMO_AUDIO):
                if not wav.exists():
                    continue
                try:
                    slide.shapes.add_movie(
                        str(wav),
                        Inches(0.6),
                        Inches(2.0 + k * 1.3),
                        Inches(1.0),
                        Inches(1.0),
                        mime_type='audio/x-wav',
                    )
                    audio_embedded += 1
                except Exception as exc:
                    print(f'  audio embed failed ({wav.name}): {exc}')

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    mb = OUT.stat().st_size / 1e6
    print(
        f'Wrote {OUT} ({mb:.1f} MB) — {len(shots)} slides, audio clips embedded: {audio_embedded}'
    )


if __name__ == '__main__':
    main()
